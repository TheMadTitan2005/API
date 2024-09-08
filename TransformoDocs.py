import json
import os
import base64
import magic
import easyocr
import fitz
import numpy as np
from PIL import Image
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation

class TransformoDocs:
    def __init__(self):
        self.magic = magic.Magic()
        self.reader = easyocr.Reader(['en', 'hi'])

    def process_file(self, file_path):
        if self.is_machine_readable(file_path):
            return {"error": f"The file '{os.path.basename(file_path)}' is already machine-readable."}

        content = {"text": [], "images": [], "tables": []}
        file_format = self.identify_file_format(file_path)

        if "image" in file_format.lower():
            content.update(self.process_image(file_path))
        elif "text" in file_format.lower():
            content.update(self.process_text(file_path))
        elif "word" in file_format.lower():
            content.update(self.process_docx(file_path))
        elif "pdf" in file_format.lower():
            content.update(self.process_pdf(file_path))
        elif "powerpoint" in file_format.lower():
            content.update(self.process_pptx(file_path))
        else:
            return {"error": "Unsupported file type."}

        # Ensure 'text' is a list of strings
        if isinstance(content.get("text"), str):
            content["text"] = [content["text"]]

        return content


    def is_machine_readable(self, file_path):
        machine_readable_extensions = ['.c', '.cpp', '.py', '.js', '.java', '.go', '.rb', '.php', '.swift', 'json']
        return any(file_path.endswith(ext) for ext in machine_readable_extensions)

    def identify_file_format(self, file_path):
        try:
            mime_type = self.magic.from_file(file_path)
            return mime_type
        except Exception as e:
            return None

    def process_image(self, file_path):
        with open(file_path, "rb") as img_file:
            image_data = img_file.read()
            image_base64 = base64.b64encode(image_data).decode('utf-8')
            ocr_result = self.reader.readtext(image_data)
            return {
                "text": [], 
                "images": [{
                    "base64": image_base64,
                    "ocr_text": " ".join([text[1] for text in ocr_result])
                }],
                "tables": [] 
            }

    def process_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            return {
                "text": [file.read()],
                "images": [], 
                "tables": [] 
            }

    def process_docx(self, docx_path):
        doc = Document(docx_path)
        content = {"text": [], "images": [], "tables": []}
        for para in doc.paragraphs:
            content["text"].append(para.text.strip())
        for table in doc.tables:
            table_data = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            content["tables"].append(table_data)
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image = rel.target_part.blob
                image_base64 = base64.b64encode(image).decode('utf-8')
                extracted_text = self.reader.readtext(image)
                ocr_text = " ".join([text[1] for text in extracted_text])
                content["images"].append({"base64": image_base64, "ocr_text": ocr_text})
        return content

    def process_pdf(self, pdf_path):
        result = {"text": [], "images": [], "tables": []}
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text.strip():
                result["text"].append(text.strip())
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                ocr_text = " ".join([text[1] for text in self.reader.readtext(image_bytes)])
                result["images"].append({
                    "page": page_num + 1,
                    "image_index": img_index + 1,
                    "base64": image_base64,
                    "ocr_text": ocr_text,
                })
        return result

    def process_pptx(self, pptx_path):
        presentation = Presentation(pptx_path)
        content = {"text": [], "images": [], "tables": []}
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    content["text"].append(shape.text_frame.text.strip())
                if shape.has_table:
                    table_data = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                    content["tables"].append(table_data)
                if shape.shape_type == 13: 
                    image = shape.image.blob
                    image_base64 = base64.b64encode(image).decode('utf-8')
                    extracted_text = self.reader.readtext(image)
                    ocr_text = " ".join([text[1] for text in extracted_text])
                    content["images"].append({"base64": image_base64, "ocr_text": ocr_text})
        return content
